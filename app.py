import os
import subprocess
import tempfile
import uuid
import logging
from fastapi import FastAPI, Request
from fastapi.responses import JSONResponse, PlainTextResponse
from pptx import Presentation
from PIL import Image
import pytesseract
from io import BytesIO
import requests
from urllib.parse import unquote, urlparse
from bs4 import BeautifulSoup

logger = logging.getLogger("uvicorn.error")
app = FastAPI()

# ---------------- PPTX Extraction ---------------- #
def extract_text_from_shape(shape):
    text = ''
    if hasattr(shape, 'text_frame') and shape.text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text += run.text
            text += '\n'
    elif shape.shape_type == 6:  # Group Shape
        for subshape in shape.shapes:
            text += extract_text_from_shape(subshape)
    elif shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                text += cell.text + '\t'
            text += '\n'
    return text

def extract_text_and_ocr_from_pptx(file_bytes):
    prs = Presentation(BytesIO(file_bytes))
    slide_texts = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_text = f"Slide {i}:\n"
        for shape in slide.shapes:
            if shape.has_text_frame or shape.has_table or shape.shape_type == 6:
                slide_text += extract_text_from_shape(shape)

            if shape.shape_type == 13:  # Picture
                img = Image.open(BytesIO(shape.image.blob))
                ocr_text = pytesseract.image_to_string(img)
                slide_text += "[OCR]: " + ocr_text + '\n'
        slide_texts.append(slide_text.strip())

    return "\n\n---\n\n".join(slide_texts)

@app.post("/extract-pptx/")
async def extract_pptx(request: Request):
    try:
        data = await request.json()
        document_url = data.get("documentUrl")
        if not document_url:
            return PlainTextResponse("Missing documentUrl", status_code=400)

        parsed_path = unquote(urlparse(document_url).path).lower()
        if not parsed_path.endswith(".pptx"):
            return PlainTextResponse("Only .pptx files are supported", status_code=400)

        response = requests.get(document_url, headers={"User-Agent": "Mozilla/5.0"})
        if response.status_code != 200:
            return PlainTextResponse(f"Failed to download file: {response.status_code}", status_code=400)

        content = extract_text_and_ocr_from_pptx(response.content)
        return PlainTextResponse(content, status_code=200)

    except Exception as e:
        logger.exception("Error during PPTX processing:")
        return PlainTextResponse(f"Processing error: {str(e)}", status_code=500)

# ---------------- Web Page Scraping ---------------- #
def fetch_and_clean_content(url):
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, 'html.parser')
        for script in soup(["script", "style"]):
            script.extract()
        lines = (line.strip() for line in soup.get_text().splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        return '\n'.join(chunk for chunk in chunks if chunk)
    except requests.exceptions.RequestException as e:
        logger.error(f"Scraping error: {e}")
        return None

@app.post("/extract-url/")
async def extract_webpage_text(request: Request):
    try:
        data = await request.json()
        url = data.get("url")
        if not url:
            return JSONResponse({"error": "Missing 'url' in request body"}, status_code=400)

        content = fetch_and_clean_content(url)
        if content:
            return JSONResponse({"url": url, "content": content}, status_code=200)
        else:
            return JSONResponse({"error": "Failed to retrieve content"}, status_code=500)
    except Exception as e:
        logger.exception("Error during URL scraping:")
        return JSONResponse({"error": str(e)}, status_code=500)

# ---------------- Dynamic Python Code Execution ---------------- #
@app.post("/execute-python/")
async def execute_python(request: Request):
    try:
        data = await request.json()
        python_code = data.get("code")
        if not python_code:
            return JSONResponse({"error": "Missing 'code' in request body"}, status_code=400)

        # Save to temporary file
        tmp_file_path = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4()}.py")
        with open(tmp_file_path, "w", encoding="utf-8") as f:
            f.write(python_code)

        # Run python code with timeout
        result = subprocess.run(
            ["python", tmp_file_path],
            capture_output=True,
            text=True,
            timeout=10
        )

        os.remove(tmp_file_path)

        if result.stderr:
            logger.error(f"Python stderr: {result.stderr}")
        output = result.stdout.strip() or "Python script did not return any output"

        return JSONResponse({"answers": [output]})

    except subprocess.TimeoutExpired:
        return JSONResponse({"answers": ["Script execution timed out"]}, status_code=408)
    except Exception as e:
        logger.exception("Python execution error:")
        return JSONResponse({"answers": [f"Script execution failed: {str(e)}"]}, status_code=500)
